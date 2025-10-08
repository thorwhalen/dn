"""
Converting things to markdown

This module provides tools to convert various file formats (PDF, Word, Excel, PowerPoint, HTML)
to markdown format. It includes:
- Default converters for common file formats
- Functions to convert bytes to markdown
- Store-based converters for processing multiple files
"""

# Note: This module is also vendored in dn package
# TODO: Keep in sync and centralize when makes sense

import json
from functools import partial
import contextlib
import base64
import io
from typing import (
    Callable,
    Dict,
    Optional,
    Mapping,
    MutableMapping,
    Any,
    Union,
    Iterable,
)

from dol import Files, TextFiles

from dn.util import (
    fullpath,
    identity,
    is_url,
    url_to_contents,
    save_to_file_and_return_file,
)

# Import libraries with error suppression
ignore_import_errors = contextlib.suppress(ImportError)

# Initialize the default converters dictionary
dflt_converters: Dict[str, Callable[[bytes], str]] = {}

dflt_md_inner_file_header = "###"

# PDF Conversion
with ignore_import_errors:
    import pypdf

    def pdf_to_markdown(
        pdf_bytes: bytes, *, md_inner_file_header=dflt_md_inner_file_header
    ) -> str:
        """Convert PDF to markdown text."""
        pdf_reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        pages = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            pages.append(f"{md_inner_file_header} Page {len(pages) + 1}\n\n{text}")
        return "\n\n".join(pages)

    dflt_converters["pdf"] = pdf_to_markdown

# Microsoft Office Conversion
with ignore_import_errors:
    import mammoth  # for .docx  # pip install mammoth

    def docx_to_markdown(docx_bytes: bytes) -> str:
        """Convert DOCX to markdown."""
        result = mammoth.convert_to_markdown(io.BytesIO(docx_bytes))
        return result.value

    dflt_converters["docx"] = docx_to_markdown
    dflt_converters["doc"] = docx_to_markdown

# Excel Conversion
with ignore_import_errors:
    import pandas as pd  # pip install pandas, openpyxl

    def excel_to_markdown(
        excel_bytes: bytes, *, md_inner_file_header=dflt_md_inner_file_header
    ) -> str:
        """Convert Excel files to markdown tables."""
        # Support both xls and xlsx
        dfs = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None)
        markdown_output = []

        for sheet_name, df in dfs.items():
            markdown_output.append(f"{md_inner_file_header} Sheet: {sheet_name}\n")
            # Convert DataFrame to markdown table
            markdown_output.append(df.to_markdown(index=False))
            markdown_output.append("\n")

        return "\n".join(markdown_output)

    dflt_converters["xlsx"] = excel_to_markdown
    dflt_converters["xls"] = excel_to_markdown

# PowerPoint Conversion
with ignore_import_errors:
    import pptx  # pip install python-pptx

    def pptx_to_markdown(
        pptx_bytes: bytes, *, md_inner_file_header=dflt_md_inner_file_header
    ) -> str:
        """Convert PowerPoint to markdown."""
        presentation = pptx.Presentation(io.BytesIO(pptx_bytes))
        slides = []

        for i, slide in enumerate(presentation.slides, 1):
            slide_text = [f"{md_inner_file_header} Slide {i}\n"]

            # Extract text from various elements
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slides.append("\n".join(slide_text))

        return "\n\n".join(slides)

    dflt_converters["pptx"] = pptx_to_markdown
    dflt_converters["ppt"] = pptx_to_markdown

# HTML Conversion
with ignore_import_errors:
    import html2text  # pip install html2text

    def html_to_markdown(html_bytes: bytes) -> str:
        """Convert HTML to markdown."""
        # Decode bytes to string
        html_str = html_bytes.decode("utf-8", errors="ignore")

        # Create HTML to Markdown converter
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False

        return h.handle(html_str)

    dflt_converters["html"] = html_to_markdown


# Fallback converter (default)
def default_fallback(b: bytes, input_format: str) -> str:
    """Fallback converter that attempts basic text extraction."""
    try:
        # Try decoding with different encodings
        encodings = ["utf-8", "latin-1", "ascii", "utf-16"]
        for encoding in encodings:
            try:
                return b.decode(encoding)
            except UnicodeDecodeError:
                continue

        # If decoding fails, convert to base64 for preservation
        return f"# Conversion Failed for {input_format}\n\n```\n{base64.b64encode(b).decode('utf-8')}\n```"
    except Exception as e:
        return f"# Conversion Error\n\nCould not convert {input_format}: {str(e)}"


# def bytes_to_markdown(
#     b: bytes,
#     input_format: str,
#     *,
#     converters: Dict[str, Callable[[bytes], str]] = dflt_converters,
#     fallback: Callable = default_fallback,
# ) -> str:
#     """
#     Convert bytes of a given format to markdown text.

#     Args:
#         b (bytes): Input bytes to convert
#         input_format (str): Format of the input (e.g., 'pdf', 'docx')
#         converters (Dict[str, Callable]): Dictionary of format-specific converters
#         fallback (Callable): Fallback conversion method

#     Returns:
#         str: Markdown-formatted text

#     See also: https://github.com/thorwhalen/aix/discussions/3#discussioncomment-12387852
#     """
#     converter = converters.get(input_format.lower(), None)
#     if converter is not None:
#         return converter(b)
#     else:
#         return fallback(b, input_format)


def add_dflt_converter(input_format: str, converter: Callable[[bytes], str]):
    """Add (or change) a default converter for a given input format."""
    dflt_converters[input_format.lower()] = converter


# --------------------------------------------------------------------------------------
# Convert Jupyter notebooks to Markdown

import os
from typing import Sequence, Callable
from contextlib import suppress

from dol import Pipe
import lkj

ignore_import_errors = suppress(ImportError, ModuleNotFoundError)

installed_packages = set()

with ignore_import_errors:
    from nbconvert.preprocessors import Preprocessor
    from nbconvert import MarkdownExporter

    installed_packages.add("nbconvert")

with ignore_import_errors:
    import nbformat

    installed_packages.add("nbformat")


def truncate_text(text: str, *, max_chars=200) -> str:
    if len(text) > max_chars:
        cutoff = max_chars
        first_part = text[: int(0.75 * cutoff)]
        last_part = text[-int(0.25 * cutoff) :]
        return f"{first_part}\n...\n{last_part}"
    return text


# TODO: Here's the default is to save md in file. Good default? What are the behaviors of dflt_converters?
# TODO: There's a hubcap.tools.notebook_to_markdown
#       --> Should these be merged?
# TODO: Also there's test2doc.notebook_utils.ensure_notebook_dict that givs us the json
# TODO: Generalize to get more control over output (and source) transformation control
def notebook_to_markdown(
    src_notebook: Optional[str] = None,
    *,
    target_file: Optional[str] = None,  # "*.md"
    output_processors: Sequence[Callable[[str], str]] = (truncate_text,),
    read_encoding: str = "utf-8",
    write_encoding: str = "utf-8",
):
    """
    Converts a Jupyter notebook to Markdown, applying a output text processors.

    Args:
        src_notebook (str): Path to the source notebook or URL.
        target_file (str, optional): Path to save the converted Markdown file. If None, returns the Markdown content.
            If it starts with '*', the '*' is replaced with the source notebook name.
            Popular value for target_file is '*.md'.
        output_processors (Sequence[Callable], optional): List of functions to process cell outputs.
        read_encoding (str): Encoding for reading the source notebook.
        write_encoding (str): Encoding for writing the Markdown file.


    """

    if src_notebook is None:
        import ipynbname  # pip install ipynbname

        src_notebook = str(ipynbname.path())

    if isinstance(src_notebook, dict):
        src_notebook = json.dumps(src_notebook)

    if isinstance(src_notebook, str) and src_notebook.startswith("~"):
        # Expand user directory
        src_notebook = os.path.expanduser(src_notebook)

    if is_url(src_notebook):
        src_content = url_to_contents(src_notebook)
        src_notebook = save_to_file_and_return_file(src_content)
    elif not isinstance(src_notebook, str) or not os.path.exists(src_notebook):
        src_content = src_notebook
        src_notebook = save_to_file_and_return_file(src_content)

    class TruncateOutputPreprocessor(Preprocessor):
        """Preprocessor that truncates long outputs in Jupyter Notebook cells."""

        def __init__(self, output_processors, **kwargs):
            super().__init__(**kwargs)
            self.output_processors = output_processors
            if self.output_processors:
                self.process_output = Pipe(*output_processors)
            else:
                self.process_output = identity

        def preprocess_cell(self, cell, resources, index):
            if "outputs" in cell:
                for output in cell.outputs:
                    # Case 1: Traditional text output
                    if isinstance(output, dict):
                        if "text" in output:
                            output["text"] = self.process_output(output["text"])

                        # Case 2: Rich MIME-based data (e.g. text/html, text/plain, etc.)
                        if "data" in output:
                            # Clear all the existing fields
                            output.clear()

                            # Convert it to a 'stream'-type output
                            # 'text' is valid only when 'output_type' is 'stream' or 'error'
                            output["output_type"] = "stream"
                            output["name"] = "stdout"  # required for 'stream'
                            output["text"] = "HTML output truncated. (Data removed)\n"
            return cell, resources

    # Load notebook
    with open(src_notebook, "r", encoding=read_encoding) as f:
        nb = nbformat.read(f, as_version=4)

    # Configure exporter with preprocessor
    exporter = MarkdownExporter()

    exporter.register_preprocessor(
        TruncateOutputPreprocessor(output_processors), enabled=True
    )

    # Convert notebook
    body, _ = exporter.from_notebook_node(nb)

    if target_file is not None:
        if target_file.startswith("*"):
            src_notebook_no_ext = os.path.splitext(src_notebook)[0]
            target_file = target_file.replace("*", src_notebook_no_ext, 1)

        # Save markdown file
        with open(target_file, "w", encoding=write_encoding) as f:
            f.write(body)

        return target_file
    else:
        # If no target file is specified, return the markdown content
        return body


dflt_converters["ipynb"] = notebook_to_markdown


# --------------------------------------------------------------------------------------

"""
Functions for detecting content types and converting various formats to markdown.
"""

import io
import binascii
import mimetypes
from typing import Optional, Callable, Dict, Any, Union
from pathlib import Path
from functools import partial


def _is_pdf(data: bytes) -> bool:
    """Check if bytes represent a PDF file by looking for the PDF header signature."""
    return data[:5] == b"%PDF-"


def _is_docx(data: bytes) -> bool:
    """Check if bytes represent a DOCX file by looking for the ZIP signature and checking if it contains word/document.xml."""
    if data[:4] != b"PK\x03\x04":  # ZIP file signature
        return False

    try:
        with io.BytesIO(data) as f:
            import zipfile

            with zipfile.ZipFile(f) as zip_ref:
                return any("word/document.xml" in name for name in zip_ref.namelist())
    except zipfile.BadZipFile:
        return False


def _is_xlsx(data: bytes) -> bool:
    """Check if bytes represent an XLSX file by looking for the ZIP signature and checking if it contains xl/workbook.xml."""
    if data[:4] != b"PK\x03\x04":  # ZIP file signature
        return False

    try:
        with io.BytesIO(data) as f:
            import zipfile

            with zipfile.ZipFile(f) as zip_ref:
                return any("xl/workbook.xml" in name for name in zip_ref.namelist())
    except zipfile.BadZipFile:
        return False


def _is_pptx(data: bytes) -> bool:
    """Check if bytes represent a PPTX file by looking for the ZIP signature and checking if it contains ppt/presentation.xml."""
    if data[:4] != b"PK\x03\x04":  # ZIP file signature
        return False

    try:
        with io.BytesIO(data) as f:
            import zipfile

            with zipfile.ZipFile(f) as zip_ref:
                return any(
                    "ppt/presentation.xml" in name for name in zip_ref.namelist()
                )
    except zipfile.BadZipFile:
        return False


def _is_html(data: bytes) -> bool:
    """Check if bytes likely represent an HTML file."""
    try:
        text = data.decode('utf-8', errors='ignore').lower()
        return (
            text.strip().startswith('<!doctype html')
            or text.strip().startswith('<html')
            or '<html' in text[:1000]
        )
    except Exception:
        return False


def _is_ipynb(data: bytes) -> bool:
    """Check if bytes represent a Jupyter notebook by checking for JSON with notebook structure."""
    try:
        import json

        notebook = json.loads(data)
        return "cells" in notebook and "metadata" in notebook and "nbformat" in notebook
    except Exception:
        return False


def _detect_with_python_magic(data: bytes) -> Optional[str]:
    """
    Attempt to detect content type using python-magic if available.

    Returns:
        String format identifier or None if detection failed or library not available
    """
    try:
        import magic

        mime = magic.from_buffer(data, mime=True)

        # Map mime types to format identifiers
        if mime == 'application/pdf':
            return 'pdf'
        elif mime in (
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/msword',
        ):
            return 'docx'
        elif mime in (
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.ms-excel',
        ):
            return 'xlsx'
        elif mime in (
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint',
        ):
            return 'pptx'
        elif mime in ('text/html', 'application/xhtml+xml'):
            return 'html'
        elif mime == 'application/json':
            # Need to check if it's specifically a Jupyter notebook
            try:
                if _is_ipynb(data):
                    return 'ipynb'
            except Exception:
                pass
            return None
    except (ImportError, Exception):
        return None

    return None


def _detect_content_type(
    data: bytes, key: Optional[str] = None, *, verbose: bool = False
) -> str:
    """
    Detect the content type of the given bytes using multiple strategies.

    Args:
        data: The byte data to analyze
        key: Optional filename that can be used for extension-based detection
        verbose: Whether to print detection information

    Returns:
        String representing the detected format (e.g., 'pdf', 'docx', etc.)
    """

    def _log(msg):
        if verbose:
            key_info = f"[{key}] " if key else ""
            print(f"{key_info}{msg}")

    # Try python-magic first if available
    magic_result = _detect_with_python_magic(data)
    if magic_result:
        _log(f"Detected content type via python-magic: {magic_result}")
        return magic_result

    # Fall back to signature-based detection
    if _is_pdf(data):
        _log("Detected PDF via signature")
        return 'pdf'
    elif _is_docx(data):
        _log("Detected DOCX via internal structure")
        return 'docx'
    elif _is_xlsx(data):
        _log("Detected XLSX via internal structure")
        return 'xlsx'
    elif _is_pptx(data):
        _log("Detected PPTX via internal structure")
        return 'pptx'
    elif _is_ipynb(data):
        _log("Detected Jupyter notebook via structure")
        return 'ipynb'
    elif _is_html(data):
        _log("Detected HTML via content pattern")
        return 'html'

    # Try to use extension from the key if available
    if key:
        extension = Path(key).suffix.lstrip('.').lower()
        if extension in (
            'pdf',
            'docx',
            'doc',
            'xlsx',
            'xls',
            'pptx',
            'ppt',
            'html',
            'ipynb',
        ):
            _log(f"Detected content type via filename extension: {extension}")
            return extension

        # Also try mimetypes module for extension-based detection
        mime_type, _ = mimetypes.guess_type(key)
        if mime_type:
            if mime_type == 'application/pdf':
                _log("Detected PDF via mimetypes")
                return 'pdf'
            elif mime_type in (
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/msword',
            ):
                _log("Detected Word document via mimetypes")
                return 'docx'
            elif mime_type in (
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'application/vnd.ms-excel',
            ):
                _log("Detected Excel document via mimetypes")
                return 'xlsx'
            elif mime_type in (
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/vnd.ms-powerpoint',
            ):
                _log("Detected PowerPoint document via mimetypes")
                return 'pptx'
            elif mime_type in ('text/html', 'application/xhtml+xml'):
                _log("Detected HTML via mimetypes")
                return 'html'

    # Additional text-based detection for JSON/Jupyter notebooks
    try:
        # Try to decode first ~1000 bytes to see if it's text
        sample = data[:1000].decode('utf-8', errors='ignore')

        # Check for JSON structure that might be a notebook
        if sample.strip().startswith('{') and '}' in sample:
            try:
                if _is_ipynb(data):
                    _log("Detected Jupyter notebook via JSON pattern")
                    return 'ipynb'
            except Exception:
                pass
    except Exception:
        pass

    _log("Could not determine content type")
    return 'unknown'


def try_to_convert_to_markdown(
    data: bytes,
    key: Optional[str] = None,
    *,
    converters: Optional[Dict[str, Callable]] = None,
    verbose: bool = False,
) -> Optional[str]:
    """
    Attempts to identify the content type of the given bytes and convert it to markdown
    using appropriate converters from dn.src.

    Args:
        data: The byte data to convert
        key: Optional identifier/filename to help with content detection and for verbose output
        converters: Optional dictionary of content type to converter function mappings
                   If None, uses the default converters from contexts module
        verbose: Whether to print information about the conversion process

    Returns:
        Markdown string if conversion successful, None otherwise


    """
    if converters is None:
        converters = dflt_converters

    def _log(msg):
        if verbose:
            key_info = f"[{key}] " if key else ""
            print(f"{key_info}{msg}")

    # First try to detect the content type
    content_type = _detect_content_type(data, key, verbose=verbose)

    if content_type == 'unknown':
        _log("Could not detect content type. Trying common formats...")

        # Try the most common formats in a reasonable order
        formats_to_try = ['pdf', 'docx', 'html', 'xlsx', 'ipynb', 'pptx']

        for fmt in formats_to_try:
            if fmt in converters:
                _log(f"Attempting conversion as {fmt}...")
                try:
                    result = bytes_to_markdown(data, fmt, converters=converters)
                    _log(f"Successfully converted content as {fmt}")
                    return result
                except Exception as e:
                    _log(f"Failed to convert as {fmt}: {str(e)}")

        _log("Could not convert content with any available converter")
        return None

    # Content type detected, try the corresponding converter
    if content_type in converters:
        try:
            result = bytes_to_markdown(data, content_type, converters=converters)
            _log(f"Successfully converted {content_type} to markdown")
            return result
        except Exception as e:
            _log(f"Failed to convert {content_type} to markdown: {str(e)}")
            return None
    else:
        _log(f"No converter available for {content_type}")
        return None


# --------------------------------------------------------------------------------------
# Markdown stores


def get_extension(path: str) -> str:
    """Get the extension of a file path."""
    return Path(path).suffix.lstrip(".").lower()


def extensions_not_supported_by_converters(
    paths: Union[str, Iterable[str]], converters: Iterable = dflt_converters
):
    """
    Returns a set of file extensions that are not supported by the given converters.

    Args:
        paths (str or Iterable): A list of file paths or the root directory to check for paths.
        converters (dict): An iterable of extensions supported by the converters.
            By default, it's the default converters (a dict, so keys are considered).

    Returns:
        set: A set of file extensions that are not supported by the given converters.

    """
    if isinstance(paths, str):
        paths = Files(paths)

    paths_extensions = {get_extension(path) for path in paths}
    supported_extensions = set(converters)

    return paths_extensions - supported_extensions


def _resolve_src_bytes_store_and_target_text_store(src_files, target_store):
    """
    Resolves the source bytes store and target text store.
    """
    if isinstance(src_files, str):
        src_rootdir = fullpath(src_files)
        src_files = Files(src_rootdir)
    else:
        if hasattr(src_files, "rootdir"):
            src_rootdir = src_files.rootdir
        else:
            src_rootdir = None

    if target_store is None:
        target_store = dict()

    if isinstance(target_store, str):
        target_store = TextFiles(target_store)

    return src_files, target_store


# def bytes_store_to_markdown_store(
#     src_files: Union[str, Mapping[str, bytes]],
#     target_store: Union[str, MutableMapping[str, bytes]] = None,
#     *,
#     converters: Dict[str, Callable[[bytes], str]] = dflt_converters,
#     verbose: bool = False,
#     old_to_new_key: Callable[[str], str] = lambda x: x + ".md",
#     target_store_egress: Callable[[Mapping], Any] = identity,
# ):
#     """
#     Converts files to markdown using the specified converters.

#     Tip: You can use `target_store_egress=aggregate_store` to get an aggregate string
#     as your output.
#     """
#     src_files, target_store = _resolve_src_bytes_store_and_target_text_store(
#         src_files, target_store
#     )
#     _bytes_to_markdown = partial(bytes_to_markdown, converters=converters)

#     for src_key, src_bytes in src_files.items():
#         if verbose:
#             print(f"Converting {src_key} to markdown")
#         ext = get_extension(src_key)
#         target_key = old_to_new_key(src_key)
#         target_store[target_key] = _bytes_to_markdown(src_bytes, ext)

#     return target_store_egress(target_store)


def bytes_to_markdown(
    b: bytes,
    input_format: Optional[str] = None,
    *,
    key: Optional[str] = None,
    converters: Dict[str, Callable[[bytes], str]] = dflt_converters,
    fallback: Callable = default_fallback,
    ext_to_input_format: Callable[[str], str] = lambda x: x,
    try_bytes_detection: bool = True,
    verbose: bool = False,
) -> str:
    """
    Convert bytes to markdown text using a flexible detection strategy.

    The function follows this logic to find a converter:
    1. If input_format is provided, use it directly to find a converter
    2. If input_format is None but key is provided, extract format using ext_to_input_format
    3. If try_bytes_detection is True, attempt content-type detection from bytes
    4. If no converter found through above methods, use the fallback

    Args:
        b (bytes): Input bytes to convert
        input_format (str, optional): Format of the input (e.g., 'pdf', 'docx')
        key (str, optional): Filename or identifier to help with format detection
        converters (Dict[str, Callable]): Dictionary of format-specific converters
        fallback (Callable): Fallback conversion method
        ext_to_input_format (Callable): Function to extract format from key
        try_bytes_detection (bool): Whether to try content detection from bytes
        verbose (bool): Whether to print detection information

    Returns:
        str: Markdown-formatted text

    Examples:
        # Convert with explicit format
        md = bytes_to_markdown(pdf_bytes, 'pdf')

        # Convert using filename extension
        md = bytes_to_markdown(file_bytes, key='document.docx')

        # Convert using only content-based detection
        md = bytes_to_markdown(file_bytes, input_format=None,
                              ext_to_input_format=None)

        # Convert with explicit format and disable content detection
        md = bytes_to_markdown(file_bytes, 'xlsx', try_bytes_detection=False)
    """

    def _log(msg):
        if verbose:
            key_info = f"[{key}] " if key else ""
            print(f"{key_info}{msg}")

    # Strategy 1: Use provided input_format
    if input_format is not None:
        _log(f"Using provided input format: {input_format}")
        format_key = input_format.lower()
        converter = converters.get(format_key)
        if converter is not None:
            return converter(b)

    # Strategy 2: Extract format from key
    if input_format is None and key is not None and ext_to_input_format is not None:
        try:
            ext = get_extension(key)
            format_key = ext_to_input_format(ext).lower()
            _log(f"Detected format from key: {format_key}")
            converter = converters.get(format_key)
            if converter is not None:
                return converter(b)
        except Exception as e:
            _log(f"Error extracting format from key: {str(e)}")

    # Strategy 3: Detect from bytes content
    if try_bytes_detection:
        try:
            _log("Attempting content-based format detection")
            content_type = _detect_content_type(b, key, verbose=verbose)
            if content_type != 'unknown':
                _log(f"Detected content type: {content_type}")
                converter = converters.get(content_type.lower())
                if converter is not None:
                    return converter(b)

            # If unknown, try common formats in a reasonable order
            if content_type == 'unknown':
                _log("Could not detect content type. Trying common formats...")
                formats_to_try = ['pdf', 'docx', 'html', 'xlsx', 'ipynb', 'pptx']

                for fmt in formats_to_try:
                    if fmt in converters:
                        _log(f"Attempting conversion as {fmt}...")
                        try:
                            result = converters[fmt](b)
                            _log(f"Successfully converted content as {fmt}")
                            return result
                        except Exception as e:
                            _log(f"Failed to convert as {fmt}: {str(e)}")
        except Exception as e:
            _log(f"Error during content detection: {str(e)}")

    # Strategy 4: Use fallback
    _log("Using fallback converter")
    detected_format = input_format or (key and get_extension(key)) or "unknown"
    return fallback(b, detected_format)


def bytes_store_to_markdown_store(
    src_files: Union[str, Mapping[str, bytes]],
    target_store: Union[str, MutableMapping[str, bytes]] = None,
    *,
    converters: Dict[str, Callable[[bytes], str]] = dflt_converters,
    fallback: Callable = default_fallback,
    ext_to_input_format: Callable[[str], str] = lambda x: x,
    try_bytes_detection: bool = True,
    verbose: bool = False,
    old_to_new_key: Callable[[str], str] = lambda x: x + ".md",
    target_store_egress: Callable[[Mapping], Any] = identity,
) -> Any:
    """
    Converts files to markdown using the enhanced bytes_to_markdown function.

    Args:
        src_files: Source files as a directory path or mapping
        target_store: Target store for markdown output
        converters: Dictionary of format-specific converters
        fallback: Fallback conversion method
        ext_to_input_format: Function to extract format from key
        try_bytes_detection: Whether to try content detection from bytes
        verbose: Whether to print detection information
        old_to_new_key: Function to transform source keys to target keys
        target_store_egress: Function to process the target store before returning

    Returns:
        Result of applying target_store_egress to the target store

    Examples:
        # Convert all files in a directory using extension detection
        result = bytes_store_to_markdown_store('path/to/files', 'path/to/output')

        # Convert using content-based detection
        result = bytes_store_to_markdown_store('path/to/files', 'path/to/output',
                                             ext_to_input_format=None)
    """
    src_files, target_store = _resolve_src_bytes_store_and_target_text_store(
        src_files, target_store
    )

    for src_key, src_bytes in src_files.items():
        if verbose:
            print(f"Converting {src_key} to markdown")

        target_key = old_to_new_key(src_key)
        target_store[target_key] = bytes_to_markdown(
            src_bytes,
            input_format=None,  # We'll detect from key or bytes
            key=src_key,
            converters=converters,
            fallback=fallback,
            ext_to_input_format=ext_to_input_format,
            try_bytes_detection=try_bytes_detection,
            verbose=verbose,
        )

    return target_store_egress(target_store)
